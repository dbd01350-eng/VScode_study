from pptx import Presentation
from pptx.util import Inches

pre = Presentation()  ##파워포인트 객체선언

for i in range(0, 11):
    title_slide_layout = pre.slide_layouts[i] # 슬라이드 종류 선택
    slide=pre.slides.add_slide(title_slide_layout)
pre.save('kim_mango.pptx')

