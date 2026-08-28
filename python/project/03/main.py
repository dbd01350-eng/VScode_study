# 라이브러리를 임포트 하고 객체를 선언
from pptx import Presentation
from pptx.util import Inches  # 사진, 표등을 그리기위해

pre = Presentation()  # 파워포인트 객체선언

for i in range(0, 11):
    title_slide_layout = pre.slide_layouts [i]
    slide=pre.slides.add_slide(title_slide_layout)

pre.save('kor.pptx')








# 제목슬라이드추가
# title=slide.placeholders
# 슬라이드 추가
# title_shape=
# 제목-제목에 값넣기
