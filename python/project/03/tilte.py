# 라이브러리를 임포트 하고 객체를 선언
from pptx import Presentation
from pptx.util import Inches  # 사진, 표등을 그리기위해

pre = Presentation()  # 파워포인트 객체선언

title_slide_layout = pre.slide_layouts[0]
slide = pre.slides.add_slide(title_slide_layout)

title = slide.placeholders[0]
title.text = "아직 한참 남았다"

subtitle = slide.placeholders[1]  # 제목상자는0 , 부제목상자는 1
subtitle.text = "와 오늘 오타 진짜 많다"

bullet_slide_layout = pre.slide_layouts[1]
slide = pre.slides.add_slide(bullet_slide_layout)

title_shape = slide.placeholders[0]
title_shape.text = "몽이는  아침에 OO를 합니다"

body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.text="몽이 아침밥먹어"

pre.save('파이썬1.pptx')
